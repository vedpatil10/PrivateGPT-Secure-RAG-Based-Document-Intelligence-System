"""
Format-specific document loaders — factory pattern supporting PDF, Word,
Excel, PowerPoint, CSV, images (OCR), emails, and plain text.
"""

import logging
from pathlib import Path
from typing import List, Dict, Optional
from abc import ABC, abstractmethod

logger = logging.getLogger("privategpt.loaders")


class LoadedDocument:
    """Represents parsed content from a document."""

    def __init__(
        self,
        content: str,
        metadata: Dict,
        page_contents: Optional[List[Dict]] = None,
    ):
        self.content = content  # Full text content
        self.metadata = metadata  # filename, type, pages, etc.
        self.page_contents = page_contents or []  # Per-page/section content


class BaseLoader(ABC):
    """Abstract base loader."""

    @abstractmethod
    def load(self, file_path: str) -> LoadedDocument:
        pass

    @abstractmethod
    def supported_extensions(self) -> List[str]:
        pass


class PDFLoader(BaseLoader):
    """Load PDF documents using pypdf with pdfplumber fallback for tables."""

    def supported_extensions(self) -> List[str]:
        return [".pdf"]

    def load(self, file_path: str) -> LoadedDocument:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        page_contents = []
        full_text_parts = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                page_contents.append({
                    "page_number": i + 1,
                    "content": text,
                })
                full_text_parts.append(text)

        # Try pdfplumber for table-heavy PDFs if pypdf yields little text
        if not full_text_parts:
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text() or ""
                        text = text.strip()
                        if text:
                            page_contents.append({
                                "page_number": i + 1,
                                "content": text,
                            })
                            full_text_parts.append(text)
            except Exception as e:
                logger.warning(f"pdfplumber fallback failed: {e}")

        full_text = "\n\n".join(full_text_parts)
        logger.info(f"PDF loaded: {Path(file_path).name}, {len(reader.pages)} pages, {len(full_text)} chars")

        return LoadedDocument(
            content=full_text,
            metadata={
                "source": file_path,
                "file_type": "pdf",
                "page_count": len(reader.pages),
            },
            page_contents=page_contents,
        )


class WordLoader(BaseLoader):
    """Load Word (.docx) documents preserving paragraph structure."""

    def supported_extensions(self) -> List[str]:
        return [".docx", ".doc"]

    def load(self, file_path: str) -> LoadedDocument:
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)
        page_contents = []
        full_text_parts = []
        current_section = ""
        current_heading = "Document Start"

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Track headings for section-based citations
            if para.style.name.startswith("Heading"):
                if current_section:
                    page_contents.append({
                        "section_title": current_heading,
                        "content": current_section,
                    })
                current_heading = text
                current_section = ""
            else:
                current_section += text + "\n"
                full_text_parts.append(text)

        # Don't forget the last section
        if current_section:
            page_contents.append({
                "section_title": current_heading,
                "content": current_section,
            })

        full_text = "\n\n".join(full_text_parts)
        logger.info(f"Word loaded: {Path(file_path).name}, {len(full_text)} chars")

        return LoadedDocument(
            content=full_text,
            metadata={
                "source": file_path,
                "file_type": "docx",
                "section_count": len(page_contents),
            },
            page_contents=page_contents,
        )


class ExcelLoader(BaseLoader):
    """Load Excel files, processing each sheet as a section."""

    def supported_extensions(self) -> List[str]:
        return [".xlsx", ".xls"]

    def load(self, file_path: str) -> LoadedDocument:
        import pandas as pd

        xls = pd.ExcelFile(file_path)
        page_contents = []
        full_text_parts = []

        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            # Convert DataFrame to readable text
            text = f"Sheet: {sheet_name}\n"
            text += df.to_string(index=False, max_rows=500)
            text = text.strip()

            if text:
                page_contents.append({
                    "section_title": f"Sheet: {sheet_name}",
                    "content": text,
                })
                full_text_parts.append(text)

        full_text = "\n\n".join(full_text_parts)
        logger.info(f"Excel loaded: {Path(file_path).name}, {len(xls.sheet_names)} sheets")

        return LoadedDocument(
            content=full_text,
            metadata={
                "source": file_path,
                "file_type": "xlsx",
                "sheet_count": len(xls.sheet_names),
            },
            page_contents=page_contents,
        )


class PowerPointLoader(BaseLoader):
    """Load PowerPoint presentations slide by slide."""

    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]

    def load(self, file_path: str) -> LoadedDocument:
        from pptx import Presentation

        prs = Presentation(file_path)
        page_contents = []
        full_text_parts = []

        for i, slide in enumerate(prs.slides):
            slide_text_parts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text_parts.append(text)

            # Also extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text_parts.append(f"[Speaker Notes: {notes}]")

            slide_text = "\n".join(slide_text_parts)
            if slide_text:
                page_contents.append({
                    "page_number": i + 1,
                    "section_title": f"Slide {i + 1}",
                    "content": slide_text,
                })
                full_text_parts.append(slide_text)

        full_text = "\n\n".join(full_text_parts)
        logger.info(f"PPTX loaded: {Path(file_path).name}, {len(prs.slides)} slides")

        return LoadedDocument(
            content=full_text,
            metadata={
                "source": file_path,
                "file_type": "pptx",
                "slide_count": len(prs.slides),
            },
            page_contents=page_contents,
        )


class CSVLoader(BaseLoader):
    """Load CSV files with column-aware context."""

    def supported_extensions(self) -> List[str]:
        return [".csv"]

    def load(self, file_path: str) -> LoadedDocument:
        import pandas as pd

        df = pd.read_csv(file_path)
        # Create column-aware text representation
        text = f"CSV with {len(df)} rows and columns: {', '.join(df.columns)}\n\n"
        text += df.to_string(index=False, max_rows=1000)

        logger.info(f"CSV loaded: {Path(file_path).name}, {len(df)} rows")

        return LoadedDocument(
            content=text,
            metadata={
                "source": file_path,
                "file_type": "csv",
                "row_count": len(df),
                "columns": list(df.columns),
            },
            page_contents=[{"content": text}],
        )


class ImageLoader(BaseLoader):
    """Load images with OCR text extraction using Tesseract."""

    def supported_extensions(self) -> List[str]:
        return [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"]

    def load(self, file_path: str) -> LoadedDocument:
        try:
            from PIL import Image
            import pytesseract

            image = Image.open(file_path)
            text = pytesseract.image_to_string(image).strip()

            logger.info(f"Image OCR: {Path(file_path).name}, {len(text)} chars extracted")

            return LoadedDocument(
                content=text if text else "[No text detected in image]",
                metadata={
                    "source": file_path,
                    "file_type": "image",
                    "image_size": f"{image.width}x{image.height}",
                },
                page_contents=[{"content": text}] if text else [],
            )
        except Exception as e:
            logger.error(f"OCR failed for {file_path}: {e}")
            return LoadedDocument(
                content=f"[OCR failed: {str(e)}]",
                metadata={"source": file_path, "file_type": "image"},
            )


class EmailLoader(BaseLoader):
    """Load .msg email files extracting subject, body, and attachment info."""

    def supported_extensions(self) -> List[str]:
        return [".msg", ".eml"]

    def load(self, file_path: str) -> LoadedDocument:
        ext = Path(file_path).suffix.lower()

        if ext == ".msg":
            return self._load_msg(file_path)
        elif ext == ".eml":
            return self._load_eml(file_path)
        else:
            raise ValueError(f"Unsupported email format: {ext}")

    def _load_msg(self, file_path: str) -> LoadedDocument:
        import extract_msg

        msg = extract_msg.Message(file_path)
        parts = [
            f"From: {msg.sender}",
            f"To: {msg.to}",
            f"Subject: {msg.subject}",
            f"Date: {msg.date}",
            "",
            msg.body or "",
        ]

        text = "\n".join(parts).strip()
        attachments = [att.longFilename for att in msg.attachments] if msg.attachments else []

        logger.info(f"Email loaded: {msg.subject}, {len(attachments)} attachments")

        return LoadedDocument(
            content=text,
            metadata={
                "source": file_path,
                "file_type": "email",
                "subject": msg.subject,
                "attachments": attachments,
            },
            page_contents=[{"content": text, "section_title": msg.subject}],
        )

    def _load_eml(self, file_path: str) -> LoadedDocument:
        import email
        from email import policy

        with open(file_path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=policy.default)

        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body = part.get_content()
                    break
        else:
            body = msg.get_content()

        parts = [
            f"From: {msg['from']}",
            f"To: {msg['to']}",
            f"Subject: {msg['subject']}",
            f"Date: {msg['date']}",
            "",
            body or "",
        ]

        text = "\n".join(parts).strip()

        return LoadedDocument(
            content=text,
            metadata={
                "source": file_path,
                "file_type": "email",
                "subject": msg["subject"],
            },
            page_contents=[{"content": text, "section_title": msg["subject"]}],
        )


class TextLoader(BaseLoader):
    """Load plain text and markdown files."""

    def supported_extensions(self) -> List[str]:
        return [".txt", ".md", ".rst", ".log"]

    def load(self, file_path: str) -> LoadedDocument:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read().strip()

        logger.info(f"Text loaded: {Path(file_path).name}, {len(text)} chars")

        return LoadedDocument(
            content=text,
            metadata={
                "source": file_path,
                "file_type": Path(file_path).suffix.lstrip("."),
            },
            page_contents=[{"content": text}],
        )


# ── Loader Factory ───────────────────────────────────────────────

_LOADERS: List[BaseLoader] = [
    PDFLoader(),
    WordLoader(),
    ExcelLoader(),
    PowerPointLoader(),
    CSVLoader(),
    ImageLoader(),
    EmailLoader(),
    TextLoader(),
]

# Build extension → loader mapping
_EXT_MAP: Dict[str, BaseLoader] = {}
for loader in _LOADERS:
    for ext in loader.supported_extensions():
        _EXT_MAP[ext] = loader

SUPPORTED_EXTENSIONS = list(_EXT_MAP.keys())


def get_loader(file_path: str) -> BaseLoader:
    """Get the appropriate loader for a file based on its extension."""
    ext = Path(file_path).suffix.lower()
    loader = _EXT_MAP.get(ext)
    if not loader:
        raise ValueError(
            f"Unsupported file type: '{ext}'. "
            f"Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )
    return loader


def load_document(file_path: str) -> LoadedDocument:
    """Load a document using the appropriate format-specific loader."""
    loader = get_loader(file_path)
    return loader.load(file_path)
